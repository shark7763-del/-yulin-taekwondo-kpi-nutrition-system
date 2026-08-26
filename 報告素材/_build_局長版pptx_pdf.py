# -*- coding: utf-8 -*-
import os
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import font_manager as fm
from matplotlib.patches import FancyBboxPatch, Rectangle, Circle
from matplotlib.backends.backend_pdf import PdfPages
from pptx import Presentation
from pptx.util import Inches

OUT = r"D:\yulin-taekwondo-kpi-nutrition-system\報告素材"
TMP = os.path.dirname(os.path.abspath(__file__))

# ---- fonts ----
def load(path, fallback=None):
    return fm.FontProperties(fname=path) if os.path.exists(path) else fallback
FP  = load(r"C:\Windows\Fonts\msjh.ttc")
FPB = load(r"C:\Windows\Fonts\msjhbd.ttc", FP)

# ---- palette ----
BG="#0e1117"; CARD="#171c26"; LINE="#2a3140"
INK="#f2f5fa"; SUB="#9aa6b8"; MUTED="#6b7686"
BRAND="#e63946"; BRAND2="#ff6b6b"
GREEN="#37d67a"; YELLOW="#ffc93c"; ORANGE="#ff9f43"; RED="#ff5b5b"
ACCENT="#4dabf7"; GOLD="#ffd479"
W,H=1280,720

def slide():
    fig=plt.figure(figsize=(12.8,7.2),dpi=200)
    ax=fig.add_axes([0,0,1,1]); ax.set_xlim(0,W); ax.set_ylim(0,H); ax.axis("off")
    ax.add_patch(Rectangle((0,0),W,H,fc=BG,ec="none"))
    return fig,ax

def Y(top): return H-top
def _clean(s):
    out=[]
    for ch in s:
        o=ord(ch)
        if (0x1F000<=o<=0x1FAFF or 0x2600<=o<=0x27BF or 0x2B00<=o<=0x2BFF
                or o in (0x200D,0xFE0F,0x2713,0x2717,0x2248) or 0x1F1E6<=o<=0x1F1FF):
            continue
        out.append(ch)
    return "".join(out).lstrip(" ")
def txt(ax,x,top,s,size=18,color=INK,bold=False,ha="left",va="top",fp=None):
    ax.text(x,Y(top),_clean(s),fontsize=size,color=color,ha=ha,va=va,
            fontproperties=(fp or (FPB if bold else FP)))
def card(ax,x,top,w,h,fc=CARD,ec=LINE,lw=1.2,r=18):
    ax.add_patch(FancyBboxPatch((x,Y(top)-h),w,h,
        boxstyle=f"round,pad=0,rounding_size={r}",fc=fc,ec=ec,lw=lw,
        mutation_aspect=1))
def pill(ax,x,top,label,fc="#12202e",ec="#2f5d86",tc=ACCENT,size=15,padx=16):
    w=len(label)*size*1.05+padx*2
    ax.add_patch(FancyBboxPatch((x,Y(top)-38),w,34,boxstyle="round,pad=0,rounding_size=17",
        fc=fc,ec=ec,lw=1))
    txt(ax,x+w/2,top+7,label,size=size,color=tc,bold=True,ha="center")
    return x+w+12
def kicker(ax,s): txt(ax,90,70,s,size=17,color=BRAND2,bold=True)
def barline(ax): ax.add_patch(Rectangle((90,Y(112)),70,4,fc=BRAND2,ec="none"))

MARGIN=90

# ============ charts ============
def draw_gradient(ax,x0,top0,w,h):
    cats=["綠燈","黃燈","橘燈","紅燈"]
    pain=[0.69,2.85,5.91,8.50]; ready=[82.14,63.48,49.36,36.00]
    bx=x0; by_top=top0; plotx=x0+40; plotw=w-70
    py0=Y(top0+h-34); py1=Y(top0+10)   # bottom, top in axis coords
    def X(k): return plotx+(k+0.5)/len(cats)*plotw
    def Yp(v): return py0+v/10*(py1-py0)
    def Yr(v): return py0+v/100*(py1-py0)
    for gl in [0,2,4,6,8,10]:
        ax.plot([plotx,plotx+plotw],[Yp(gl),Yp(gl)],color=LINE,lw=1,zorder=1)
    bw=54
    for k,v in enumerate(pain):
        ax.add_patch(Rectangle((X(k)-bw/2,Yp(0)),bw,Yp(v)-Yp(0),fc=BRAND2,ec="none",alpha=.9,zorder=2))
        ax.text(X(k),Yp(v)+10,f"{v:.2f}",ha="center",va="bottom",color=BRAND2,fontsize=15,fontproperties=FPB)
    xr=[X(k) for k in range(len(cats))]; yr=[Yr(v) for v in ready]
    ax.plot(xr,yr,color=ACCENT,lw=3,zorder=3)
    for k,v in enumerate(ready):
        ax.add_patch(Circle((X(k),Yr(v)),6,fc=ACCENT,ec="none",zorder=4))
        ax.text(X(k),Yr(v)+14,f"{v:.0f}",ha="center",va="bottom",color=ACCENT,fontsize=15,fontproperties=FPB)
    for k,c in enumerate(cats):
        ax.text(X(k),Yp(0)-18,c,ha="center",va="top",color=SUB,fontsize=16,fontproperties=FPB)

def draw_case(ax,x0,top0,w,h):
    days=["6/25","6/26","6/27","6/28","6/29","6/30","7/1","7/2","7/3"]
    ready=[75,53,37,35,48,68,50,45,45]; pain=[1,6,9,8,8,5,8,7,8]; flags=[0,0,1,1,1,0,1,1,1]
    plotx=x0+44; plotw=w-70
    py0=Y(top0+h-40); py1=Y(top0+16)
    def X(k): return plotx+k/(len(days)-1)*plotw
    def Yr(v): return py0+v/100*(py1-py0)
    def Yp(v): return py0+v/10*(py1-py0)
    for gl in [0,25,50,75,100]:
        ax.plot([plotx,plotx+plotw],[Yr(gl),Yr(gl)],color=LINE,lw=1,zorder=1)
        ax.text(plotx-12,Yr(gl),str(gl),ha="right",va="center",color=MUTED,fontsize=12,fontproperties=FP)
    for k,f in enumerate(flags):
        if f: ax.add_patch(Rectangle((X(k)-18,py1),36,py0-py1,fc=RED,ec="none",alpha=.10,zorder=1))
    ax.plot([X(k) for k in range(len(days))],[Yr(v) for v in ready],color=ACCENT,lw=3.2,zorder=3)
    for k,v in enumerate(ready): ax.add_patch(Circle((X(k),Yr(v)),4.5,fc=ACCENT,ec="none",zorder=4))
    ax.plot([X(k) for k in range(len(days))],[Yp(v) for v in pain],color=BRAND2,lw=3.2,zorder=3)
    for k,v in enumerate(pain): ax.add_patch(Circle((X(k),Yp(v)),4.5,fc=BRAND2,ec="none",zorder=4))
    for k,f in enumerate(flags):
        if f: ax.scatter([X(k)],[Yp(pain[k])+22],marker="v",s=90,color=RED,zorder=5,edgecolors="none")
    for k,d in enumerate(days):
        ax.text(X(k),py0-16,d,ha="center",va="top",color=SUB,fontsize=13,fontproperties=FP)

# ============ slides ============
figs=[]

# 1 封面
def s1():
    fig,ax=slide()
    txt(ax,MARGIN,150,"認識一位基層專任教練",size=18,color=BRAND2,bold=True)
    txt(ax,MARGIN,205,"我只是想，",size=58,color=INK,bold=True)
    txt(ax,MARGIN,278,"把每一個孩子都帶好。",size=58,color=INK,bold=True)
    txt(ax,MARGIN,380,"育林國中技擊隊專任教練，為了看懂選手每天的身心狀態，自己動手做了",size=21,color=SUB)
    txt(ax,MARGIN,414,"一套 KPI 回報系統——一天一天，把對孩子的在乎記錄下來。",size=21,color=SUB)
    x=MARGIN
    for t in ["43 名選手","713 筆真實資料","信度 α=.905","受傷早期示警","已成研究論文"]:
        x=pill(ax,x,500,t)
    figs.append(fig)

# 2 動機
def s2():
    fig,ax=slide(); kicker(ax,"起心動念 · 我為什麼做這個"); barline(ax)
    txt(ax,MARGIN,175,"一開始，我不是想做系統",size=40,color=INK,bold=True)
    ax.add_patch(Rectangle((MARGIN,Y(340)),5,120,fc=BRAND,ec="none"))
    txt(ax,MARGIN+22,255,"我只是想知道——我的選手今天到底累不累、",size=27,color=INK,bold=True)
    txt(ax,MARGIN+22,298,"有沒有帶傷、狀態好不好。",size=27,color=INK,bold=True)
    txt(ax,MARGIN+22,338,"憑印象不夠，所以我自己動手，把它記下來、追蹤起來。",size=22,color=SUB)
    txt(ax,MARGIN,430,"做著做著，就變成現在這套系統，還意外做出了一份研究。",size=22,color=SUB)
    txt(ax,MARGIN,466,"它的初衷很單純：",size=22,color=SUB)
    txt(ax,MARGIN+270,466,"把我的孩子帶得更好。",size=22,color=GOLD,bold=True)
    txt(ax,MARGIN,560,"— 育林國中技擊隊專任教練　白天帶訓練，晚上假日一點一點做，實作驗證約一年",size=17,color=MUTED)
    figs.append(fig)

# 3 問題
def s3():
    fig,ax=slide(); kicker(ax,"01 · 這解決什麼問題"); barline(ax)
    txt(ax,MARGIN,165,"基層選手的身心狀態，往往沒人接得住",size=36,color=INK,bold=True)
    txt(ax,MARGIN,222,"國中選手正值發育與課業壓力並存的關鍵期，但基層「一人帶多人、缺乏運動科學",size=19,color=SUB)
    txt(ax,MARGIN,250,"支援」，選手的疲勞、疼痛、睡眠、情緒，多半只能靠教練現場印象判斷。",size=19,color=SUB)
    cards=[("🤕 傷害風險","「說沒事」卻帶傷訓練，","等到受傷才發現。"),
           ("😴 恢復與睡眠","睡眠不足、疲勞累積的","訊號分散、轉瞬即逝。"),
           ("🧠 身心狀態","情緒與動機波動，","缺乏留存與追蹤的依據。")]
    cw=340; gap=25; x0=MARGIN
    for i,(t,l1,l2) in enumerate(cards):
        x=x0+i*(cw+gap); card(ax,x,320,cw,150)
        txt(ax,x+26,352,t,size=21,color=INK,bold=True)
        txt(ax,x+26,398,l1,size=16,color=SUB); txt(ax,x+26,424,l2,size=16,color=SUB)
    txt(ax,MARGIN,520,"這套系統，就是要把這些「看不見的狀態」變成",size=21,color=SUB)
    txt(ax,MARGIN+560,520,"可記錄、可追蹤、可示警",size=21,color=GOLD,bold=True)
    txt(ax,MARGIN+560+330,520,"的資料。",size=21,color=SUB)
    figs.append(fig)

# 4 傳統 vs KPI
def s4():
    fig,ax=slide(); kicker(ax,"02 · 做法的核心差異"); barline(ax)
    txt(ax,MARGIN,165,"從「寫訓練日誌」升級為「可追蹤的數據」",size=36,color=INK,bold=True)
    txt(ax,MARGIN,222,"傳統日誌記「做了什麼」；KPI 系統記「狀態如何、變化怎樣、風險在哪」。",size=19,color=SUB)
    rows=[("回報方式","自由書寫、流水帳","結構化評分＋睡眠/疼痛/RPE/情緒"),
          ("客觀性","主觀、難以比較","量化，可跨人、跨日、跨隊比較"),
          ("趨勢預警","看不出走勢，警訊易漏","準備度紅黃綠燈＋受傷風險自動示警"),
          ("留存追蹤","分散、易遺失","雲端留存，可回溯、可做研究")]
    x0=MARGIN; c0=170; c1=430; c2=770; top=280; rh=78
    ax.add_patch(Rectangle((x0+c1-20,Y(top)-rh*4-14),300,14+rh*4,fc="#241a17",ec="none"))
    ax.add_patch(Rectangle((x0+c2-20,Y(top)-rh*4-14),300,14+rh*4,fc="#12241a",ec="none"))
    txt(ax,x0+c1,top-42,"📓 傳統訓練日誌",size=20,color="#e0a06a",bold=True)
    txt(ax,x0+c2,top-42,"📊 KPI 回報系統",size=20,color=GREEN,bold=True)
    for i,(k,a,b) in enumerate(rows):
        yy=top+i*rh
        ax.plot([x0,x0+1080],[Y(yy+rh-8),Y(yy+rh-8)],color=LINE,lw=1)
        txt(ax,x0,yy+18,k,size=18,color=SUB,bold=True)
        txt(ax,x0+c1,yy+18,a,size=17,color="#c9a58a")
        txt(ax,x0+c2,yy+18,b,size=17,color="#d7f5e3",bold=True)
    figs.append(fig)

# 5 五週成果
def s5():
    fig,ax=slide(); kicker(ax,"03 · 五週的真實運作成果"); barline(ax)
    txt(ax,MARGIN,165,"不是概念，是已經在跑的真實資料",size=36,color=INK,bold=True)
    txt(ax,MARGIN,222,"資料期間 2026/6/1–7/4，約五週，單一基層團隊（校隊＋道館）。",size=19,color=SUB)
    stats=[("43","參與選手",BRAND2),("713","每日回報筆數",BRAND2),
           ("88","自動風險示警(含 high 11)",GREEN),("361","AI 準備度評分",ACCENT)]
    cw=250; gap=22; x0=MARGIN
    for i,(n,l,c) in enumerate(stats):
        x=x0+i*(cw+gap); card(ax,x,300,cw,180)
        txt(ax,x+cw/2,388,n,size=62,color=c,bold=True,ha="center")
        txt(ax,x+cw/2,420,l,size=15,color=SUB,ha="center")
    txt(ax,MARGIN,540,"選手每天訓練後用手機回報，系統自動運算準備度、標記風險——",size=20,color=SUB)
    txt(ax,MARGIN,572,"五週累積出可分析、可驗證的資料庫。",size=20,color=INK,bold=True)
    figs.append(fig)

# 6 有訊號 + gradient
def s6():
    fig,ax=slide(); kicker(ax,"04 · 這些數據「有訊號」"); barline(ax)
    txt(ax,MARGIN,160,"經統計檢驗，資料具科學可信度",size=34,color=INK,bold=True)
    stats=[("α=.905","六大構面 KPI 信度（優良）",GOLD),
           ("ρ=−.51","疼痛越高、準備度越低",BRAND2),
           ("ρ=+.75","恢復越好、準備度越高",ACCENT)]
    cw=340; gap=25; x0=MARGIN
    for i,(n,l,c) in enumerate(stats):
        x=x0+i*(cw+gap); card(ax,x,205,cw,120)
        txt(ax,x+cw/2,258,n,size=40,color=c,bold=True,ha="center")
        txt(ax,x+cw/2,300,l,size=14,color=SUB,ha="center")
    card(ax,MARGIN,350,1100,270,fc=CARD)
    txt(ax,MARGIN+24,378,"紅黃綠燈分級 vs 實際身心指標",size=19,color=INK,bold=True)
    txt(ax,MARGIN+1076,378,"● 真實資料",size=14,color=GREEN,bold=True,ha="right")
    draw_gradient(ax,MARGIN+24,398,1052,205)
    txt(ax,MARGIN,650,"燈號越差，疼痛單調上升(0.7→8.5)、準備度單調下降(82→36)——證明系統燈號真的反映選手狀態。",
        size=16,color=GOLD,bold=True)
    figs.append(fig)

# 7 A12 案例
def s7():
    fig,ax=slide(); kicker(ax,"05 · 為什麼我在乎每一次疼痛"); barline(ax)
    txt(ax,MARGIN,160,"系統怎麼在受傷前，先舉起手",size=34,color=INK,bold=True)
    card(ax,MARGIN,200,1100,300,fc=CARD)
    txt(ax,MARGIN+24,228,"🥋 選手 A12（匿名）· 腰部疼痛期間逐日追蹤",size=19,color=INK,bold=True)
    txt(ax,MARGIN+1076,228,"● 真實資料 6/25–7/03",size=14,color=GREEN,bold=True,ha="right")
    draw_case(ax,MARGIN+24,250,1052,235)
    txt(ax,MARGIN,540,"● AI 準備度(左軸 0–100)    ● 疼痛(右軸 0–10)    ▼ = 系統發出 high「受傷風險(腰)」示警",
        size=15,color=SUB)
    txt(ax,MARGIN,588,"6/26 起腰痛竄升(疼痛 1→9)、準備度崩落(75→35)，系統連續多日標記「受傷風險」，我也因此更早介入、調整訓練。",size=16,color=SUB)
    txt(ax,MARGIN,616,"會想做這個，是因為我不想再看到任何一個孩子，帶著傷自己硬撐。",size=17,color=GOLD,bold=True)
    figs.append(fig)

# 8 價值
def s8():
    fig,ax=slide(); kicker(ax,"06 · 這帶來的價值"); barline(ax)
    txt(ax,MARGIN,165,"對選手照顧、對教練減負、對政策加分",size=36,color=INK,bold=True)
    cards=[("🥋 選手","身心與傷勢被看見，","訓練更安全、更被理解。"),
           ("🧑‍🏫 教練","把經驗變資料，","決策有依據，早一步介入。"),
           ("👨‍👩‍👧 家長","自動報告與 LINE，","親師生溝通有共同語言。"),
           ("🏫 學校/主管機關","基層選手照顧的示範，","留下可對外的成果。")]
    cw=250; gap=22; x0=MARGIN
    for i,(t,l1,l2) in enumerate(cards):
        x=x0+i*(cw+gap); card(ax,x,250,cw,150)
        txt(ax,x+22,282,t,size=19,color=INK,bold=True)
        txt(ax,x+22,326,l1,size=14,color=SUB); txt(ax,x+22,350,l2,size=14,color=SUB)
    ax.add_patch(Rectangle((MARGIN,Y(560)),5,110,fc=BRAND,ec="none"))
    txt(ax,MARGIN+22,470,"價值不在於 AI 取代教練，",size=25,color=INK,bold=True)
    txt(ax,MARGIN+22,510,"而在把選手狀態變成可記錄、可追蹤、可討論的資料。",size=25,color=INK,bold=True)
    figs.append(fig)

# 9 量化效益 大數字
def s9():
    fig,ax=slide(); kicker(ax,"07 · 這份在乎，也留下了痕跡"); barline(ax)
    txt(ax,MARGIN,165,"多出來的時間，我都還給孩子",size=36,color=INK,bold=True)
    heros=[("70","分鐘 / 天","省下的盤點時間，","全拿去多看孩子幾眼",BRAND2,"估算",GOLD),
           ("80%","↓ 時間","掌握全隊狀態","所需時間大幅下降",GOLD,"估算",GOLD),
           ("88","次示警","系統替我多守的","含 11 次高風險早期攔截",GREEN,"真實",GREEN)]
    cw=340; gap=25; x0=MARGIN
    for i,(n,u,c1,c2,col,badge,bcol) in enumerate(heros):
        x=x0+i*(cw+gap); card(ax,x,240,cw,230,r=20)
        txt(ax,x+cw/2,360,n,size=94,color=col,bold=True,ha="center",va="baseline")
        txt(ax,x+cw/2,388,u,size=20,color=SUB,bold=True,ha="center")
        txt(ax,x+cw/2,424,c1,size=15,color=SUB,ha="center")
        txt(ax,x+cw/2,448,c2,size=15,color=SUB,ha="center")
        ax.text(x+cw-52,Y(268),badge,fontsize=12,color=bcol,ha="left",va="center",fontproperties=FPB)
    txt(ax,MARGIN,520,"每週多出的 約 6 小時，我拿去多陪練、多留意幾個孩子；",size=19,color=INK,bold=True)
    txt(ax,MARGIN,552,"系統也替我多守了 88 道 我可能會漏看的安全防線。",size=19,color=INK,bold=True)
    txt(ax,MARGIN,600,"＊時間估算基準：過去逐一掌握每位選手狀態約 2 分鐘 × 約 40 人 ＝ 約 86 分鐘/日；導入後以選手自主回報＋",size=12.5,color=MUTED)
    txt(ax,MARGIN,622,"戰情室彙整約 15 分鐘/日估算(省約 70 分鐘、約 80%、每週約 6 小時)。示警 88 次為系統五週實際紀錄。",size=12.5,color=MUTED)
    figs.append(fig)

# 10 研究
def s10():
    fig,ax=slide(); kicker(ax,"08 · 從實務走到研究"); barline(ax)
    txt(ax,MARGIN,165,"基層教練，做出了研究成果",size=36,color=INK,bold=True)
    card(ax,MARGIN,230,1100,240,fc="#1b1a12",ec=GOLD,r=18)
    txt(ax,MARGIN+30,272,"《AI 輔助 KPI 回報系統應用於國中跆拳道選手",size=24,color=GOLD,bold=True)
    txt(ax,MARGIN+30,306,"訓練準備度監測之初探》",size=24,color=GOLD,bold=True)
    lines=["· 以後台真實資料（713 筆／43 人／五週）完成的研究論文初稿",
           "· 含信度分析（α=.905）、相關分析、分級比較與個案追蹤",
           "· 全程匿名化（A01–A43）、家長同意、恪守研究倫理",
           "· 可延伸投稿運動科學／運動科技／行動研究方向"]
    for i,l in enumerate(lines):
        txt(ax,MARGIN+30,352+i*28,l,size=16,color=SUB)
    txt(ax,MARGIN,520,"為了真的看懂孩子的狀態，我連信度分析、研究倫理都自己弄懂——",size=20,color=SUB)
    txt(ax,MARGIN,552,"因為我想對得起每一個，把身體交給我的孩子。",size=20,color=GOLD,bold=True)
    figs.append(fig)

# 11 結語
def s11():
    fig,ax=slide()
    txt(ax,MARGIN,120,"結語",size=17,color=BRAND2,bold=True)
    txt(ax,MARGIN,185,"我沒有想證明什麼，",size=50,color=INK,bold=True)
    txt(ax,MARGIN,250,"只是想把孩子帶好。",size=50,color=INK,bold=True)
    ax.add_patch(Rectangle((MARGIN,Y(500)),5,150,fc=BRAND,ec="none"))
    txt(ax,MARGIN+22,375,"一位基層專任教練，用約一年的時間，",size=24,color=INK,bold=True)
    txt(ax,MARGIN+22,412,"把對每一個孩子的在乎，變成看得見的紀錄。",size=24,color=INK,bold=True)
    txt(ax,MARGIN+22,449,"如果這份用心也能幫到其他基層的孩子，我很樂意分享。",size=24,color=INK,bold=True)
    txt(ax,MARGIN,555,"育林國中技擊隊　專任教練　報告人：楊復傑",size=18,color=MUTED)
    figs.append(fig)

for f in [s1,s2,s3,s4,s5,s6,s7,s8,s9,s10,s11]: f()

# ---- export PNGs ----
pngs=[]
for i,fig in enumerate(figs,1):
    p=os.path.join(TMP,f"slide_{i:02d}.png")
    fig.savefig(p,dpi=200,facecolor=BG); pngs.append(p)

# ---- PDF ----
pdf_path=os.path.join(OUT,"局長展示版簡報.pdf")
with PdfPages(pdf_path) as pdf:
    for fig in figs: pdf.savefig(fig,facecolor=BG)
print("PDF ->",pdf_path)

# ---- PPTX (full-bleed images) ----
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
blank=prs.slide_layouts[6]
for p in pngs:
    s=prs.slides.add_slide(blank)
    s.shapes.add_picture(p,0,0,width=prs.slide_width,height=prs.slide_height)
pptx_path=os.path.join(OUT,"局長展示版簡報.pptx")
prs.save(pptx_path)
print("PPTX ->",pptx_path)
print("slides:",len(figs))
